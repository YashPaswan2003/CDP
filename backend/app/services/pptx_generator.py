"""PPTX file generator for presentations using python-pptx."""

from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
import logging

logger = logging.getLogger("api")

# Brand colors
DARK_BG = RGBColor(0x0A, 0x0A, 0x0F)
SURFACE_BG = RGBColor(0x14, 0x14, 0x1E)
PRIMARY_BLUE = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)
TEXT_PRIMARY = RGBColor(0xF8, 0xFA, 0xFC)
TEXT_SECONDARY = RGBColor(0x94, 0xA3, 0xB8)
SUCCESS_GREEN = RGBColor(0x22, 0xC5, 0x5E)
DANGER_RED = RGBColor(0xEF, 0x44, 0x44)


def _set_slide_bg(slide, color=DARK_BG):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, left, top, width, height, text, font_size=14,
                  color=TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def _add_bullet_list(slide, left, top, width, height, bullets, font_size=12):
    """Add a bulleted list to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(6)
        p.level = 0
    return txBox


def _add_metric_card(slide, left, top, label, value, change=None):
    """Add a metric card (rounded rectangle with label + value)."""
    card_w = Inches(2.2)
    card_h = Inches(1.0)

    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.ROUNDED_RECTANGLE
        left, top, card_w, card_h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = SURFACE_BG
    shape.line.fill.background()

    # Value
    _add_text_box(slide, left + Inches(0.15), top + Inches(0.1),
                  card_w - Inches(0.3), Inches(0.4),
                  str(value), font_size=18, bold=True, color=TEXT_PRIMARY)

    # Label + change
    label_text = label
    if change:
        label_text += f"  ({change})"
    _add_text_box(slide, left + Inches(0.15), top + Inches(0.55),
                  card_w - Inches(0.3), Inches(0.35),
                  label_text, font_size=9, color=TEXT_SECONDARY)


def generate_title_slide(prs, client_name, date_from, date_to):
    """Generate title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide)

    # Client name
    _add_text_box(slide, Inches(1), Inches(2.0), Inches(8), Inches(1),
                  client_name, font_size=36, bold=True, color=TEXT_PRIMARY,
                  alignment=PP_ALIGN.CENTER)

    # Subtitle
    _add_text_box(slide, Inches(1), Inches(3.0), Inches(8), Inches(0.5),
                  "Performance Report", font_size=24, color=PRIMARY_BLUE,
                  alignment=PP_ALIGN.CENTER)

    # Date range
    _add_text_box(slide, Inches(1), Inches(3.8), Inches(8), Inches(0.5),
                  f"{date_from} to {date_to}", font_size=14, color=TEXT_SECONDARY,
                  alignment=PP_ALIGN.CENTER)

    # Footer
    _add_text_box(slide, Inches(1), Inches(6.5), Inches(8), Inches(0.4),
                  "Prepared by Ethinos Digital", font_size=10, color=TEXT_SECONDARY,
                  alignment=PP_ALIGN.CENTER)


def generate_metrics_slide(prs, slide_data):
    """Generate a slide with metric cards and bullets."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide)

    title = slide_data.get("title", "Metrics")
    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                  title, font_size=24, bold=True, color=TEXT_PRIMARY)

    # Metric cards row
    metrics = slide_data.get("metrics", [])
    for i, m in enumerate(metrics[:4]):
        left = Inches(0.5 + i * 2.4)
        _add_metric_card(slide, left, Inches(1.2),
                         m.get("label", ""), m.get("value", ""),
                         m.get("change", None))

    # Bullets below
    bullets = slide_data.get("bullets", [])
    if bullets:
        _add_bullet_list(slide, Inches(0.5), Inches(2.5), Inches(9), Inches(4.5),
                         bullets, font_size=13)


def generate_table_slide(prs, slide_data):
    """Generate a slide with a data table."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide)

    title = slide_data.get("title", "Data Table")
    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                  title, font_size=24, bold=True, color=TEXT_PRIMARY)

    rows_data = slide_data.get("table_rows", [])
    headers = slide_data.get("table_headers", [])

    if not rows_data or not headers:
        _add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(1),
                      "No data available", font_size=14, color=TEXT_SECONDARY)
        return

    n_rows = min(len(rows_data) + 1, 11)  # Max 10 data rows + header
    n_cols = len(headers)

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.5), Inches(1.2),
        Inches(9), Inches(min(n_rows * 0.45, 5.5))
    )
    table = table_shape.table

    # Style header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = str(header)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.bold = True
            paragraph.font.color.rgb = TEXT_PRIMARY
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1E, 0x40, 0xAF)

    # Data rows
    for i, row in enumerate(rows_data[:10]):
        for j, val in enumerate(row[:n_cols]):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.color.rgb = TEXT_PRIMARY
            cell.fill.solid()
            cell.fill.fore_color.rgb = SURFACE_BG if i % 2 == 0 else DARK_BG


def generate_bullet_slide(prs, slide_data):
    """Generate a slide with title and bullet points."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide)

    title = slide_data.get("title", "Content")
    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                  title, font_size=24, bold=True, color=TEXT_PRIMARY)

    bullets = slide_data.get("bullets", [])
    if bullets:
        _add_bullet_list(slide, Inches(0.7), Inches(1.3), Inches(8.5), Inches(5.5),
                         bullets, font_size=14)


def generate_pptx(presentation_data: dict) -> BytesIO:
    """Generate a complete PPTX file from presentation data.

    Args:
        presentation_data: Dict with title, client_name, date_from, date_to, slides[]

    Returns:
        BytesIO buffer containing the PPTX file
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    client_name = presentation_data.get("client_name", "Client")
    date_from = presentation_data.get("date_from", "")
    date_to = presentation_data.get("date_to", "")

    # Title slide
    generate_title_slide(prs, client_name, date_from, date_to)

    # Content slides
    slides = presentation_data.get("slides", [])
    for slide_data in slides:
        slide_type = slide_data.get("type", "custom")

        if slide_type in ("executive_summary", "platform_overview", "funnel_analysis"):
            generate_metrics_slide(prs, slide_data)
        elif slide_type == "campaign_performance":
            if slide_data.get("table_rows"):
                generate_table_slide(prs, slide_data)
            else:
                generate_metrics_slide(prs, slide_data)
        elif slide_type == "keyword_analysis":
            if slide_data.get("table_rows"):
                generate_table_slide(prs, slide_data)
            else:
                generate_bullet_slide(prs, slide_data)
        elif slide_type == "recommendations":
            generate_bullet_slide(prs, slide_data)
        else:
            # custom or unknown type
            if slide_data.get("table_rows"):
                generate_table_slide(prs, slide_data)
            elif slide_data.get("metrics"):
                generate_metrics_slide(prs, slide_data)
            else:
                generate_bullet_slide(prs, slide_data)

    # Save to buffer
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
